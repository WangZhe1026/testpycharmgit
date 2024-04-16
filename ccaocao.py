import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os
import re

def pdf_to_ppt(pdf_path, ppt_path):
    # 打开PDF文件
    doc = fitz.open(pdf_path)
    ppt = Presentation()

    # for page in doc:  # 遍历PDF的每一页
    #     pix = page.get_pixmap()  # 将PDF页面转换为图像
    #     img_path = f"page_{page.number}.png"
    #     pix.save(img_path)  # 保存图像文件
    #
    #     # 创建新的幻灯片并插入图像
    #     slide = ppt.slides.add_slide(ppt.slide_layouts[6])  # 使用空白幻灯片布局
    #     slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=ppt.slide_width)

    for page in doc:

        # 创建一个Matrix对象，用于调整图像的缩放
        zoom_x = 20.0  # 横向缩放的倍数
        zoom_y = 20.0  # 纵向缩放的倍数
        mat = fitz.Matrix(zoom_x, zoom_y)  # 创建缩放矩阵

        # 将PDF页面转换为图像
        pix = page.get_pixmap(matrix=mat, alpha=False)  # alpha=False表示不使用透明度

        # 生成图像文件路径
        img_path = f"page_{page.number}.png"
        # 保存图像文件
        pix.save(img_path)
        # 创建新的幻灯片并插入图像
        slide = ppt.slides.add_slide(ppt.slide_layouts[6])  # 使用空白幻灯片布局
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=ppt.slide_width)
    ppt.save(ppt_path)  # 保存PPT文件

    # 删除生成的图像文件
    # 设置文件所在的目录
    directory = os.getcwd()

    # 正则表达式匹配文件名
    pattern = re.compile(r"^page_\d+\.png$")
    for filename in os.listdir(directory):
        if pattern.match(filename):  # 检查文件名是否匹配
            file_path = os.path.join(directory, filename)  # 构造文件的完整路径
            os.remove(file_path)  # 删除文件
            print(f"Deleted {file_path}")  # 打印删除信息

# 使用示例
# woshishabi
if  name =='main':
    pdf_to_ppt("slide.pdf", "output.pptx")
